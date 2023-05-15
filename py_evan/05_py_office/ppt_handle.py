# coding:utf-8
# time: 2023/5/14
# author: evan
# python操作PPT

import pptx
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt, Inches


class PPTHandler:
    def __init__(self, path: str):
        self.path = path
        self.ppt = pptx.Presentation()

    def create_slide(self):
        """
        创建一张ppt
        :return:
        """
        layout = self.ppt.slide_layouts[0]
        slide = self.ppt.slides.add_slide(layout)
        return slide

    def add_paragraph(self):
        slide = self.create_slide()
        placeholder = slide.placeholders[0]
        placeholder.text = 'this is new content,author:evan'

        paragraph1 = placeholder.text_frame.add_paragraph()
        paragraph1.text = '欢迎学习ppt制作'
        paragraph1.bold = True
        paragraph1.font.italic = True
        paragraph1.font.size = Pt(16)
        paragraph1.font.underline = True
        paragraph1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        paragraph2 = placeholder.text_frame.add_paragraph()
        paragraph2.text = '欢迎学习python'
        paragraph2.font.size = Pt(32)
        paragraph2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    def add_table(self):
        slide = self.create_slide()
        rows = 10
        cols = 2

        left = top = Inches(2)
        width = Inches(6.0)
        height = Inches(1.0)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for index, _ in enumerate(range(rows)):
            for sub_index in range(cols):
                table.cell(index, sub_index).text = '%s:%s' % (index, sub_index)

    def add_image(self):
        slide = self.create_slide()
        image = slide.shapes.add_picture(
            image_file='./doc/logo2020.png',
            left=Inches(1),
            top=Inches(1),
            width=Inches(6),
            height=Inches(4)
        )

    def read_ppt(self):
        p = pptx.Presentation(self.path)
        for slide in p.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    print(f'text:{shape.text_frame.text}')
                if shape.has_table:
                    for cell in shape.table.iter_cells():
                        print(cell.text, end='\t')
                    print('')

    def __del__(self):
        self.ppt.save(self.path)


if __name__ == '__main__':
    ppt = PPTHandler('leon.ppt')
    ppt.create_slide()
    ppt.add_paragraph()
    ppt.add_table()
    ppt.add_image()
    ppt.read_ppt()
